import os
import re
from typing import Dict, Any, Literal
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from datetime import datetime
import subprocess


class DocumentGenerator:
    """
    Generate Word and PowerPoint documents from intelligence data
    """
    
    def __init__(self):
        # Ensure output directory exists
        os.makedirs("output", exist_ok=True)
    
    async def generate_document(
        self,
        intelligence: Dict[str, Any],
        query: str,
        research_type: Literal["individual", "organization"],
        output_format: Literal["word", "pdf", "powerpoint"]
    ) -> str:
        """
        Generate document in specified format
        """
        if output_format == "powerpoint":
            return self._generate_powerpoint(intelligence, query, research_type)
        elif output_format == "pdf":
            # Generate Word first, then convert to PDF
            word_path = self._generate_word(intelligence, query, research_type)
            pdf_path = word_path.replace(".docx", ".pdf")
            self._convert_word_to_pdf(word_path, pdf_path)
            return pdf_path
        else:
            return self._generate_word(intelligence, query, research_type)
    
    def _generate_word(
        self,
        intelligence: Dict[str, Any],
        query: str,
        research_type: Literal["individual", "organization"]
    ) -> str:
        """
        Generate Word document
        """
        doc = Document()
        
        # Set document margins
        sections = doc.sections
        for section in sections:
            section.top_margin = Inches(1)
            section.bottom_margin = Inches(1)
            section.left_margin = Inches(1)
            section.right_margin = Inches(1)
        
        # Title
        title = doc.add_heading(f"Sales Enablement Intelligence Capsule", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Subtitle
        subtitle = doc.add_heading(f"{research_type.title()}: {query}", 1)
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Date
        date_para = doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_page_break()
        
        if research_type == "individual":
            self._add_individual_content(doc, intelligence)
        else:
            self._add_organization_content(doc, intelligence)
        
        # Save document
        # Clean query for filename - remove special characters and limit length
        clean_query = re.sub(r'[^\w\s-]', '', query)  # Remove special chars
        clean_query = re.sub(r'\s+', '_', clean_query)  # Replace spaces with underscores
        clean_query = clean_query[:50]  # Limit length
        clean_query = clean_query.strip('_')  # Remove leading/trailing underscores
        
        # If query is empty after cleaning, use a default
        if not clean_query:
            clean_query = "research"
        
        filename = f"{clean_query}_{research_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        filepath = f"output/{filename}"
        
        try:
            doc.save(filepath)
        except Exception as e:
            print(f"❌ Error saving document: {e}")
            import traceback
            print(f"   Traceback: {traceback.format_exc()}")
            raise
        
        return filepath
    
    def _add_individual_content(self, doc: Document, intelligence: Dict[str, Any]):
        """
        Add individual research content to Word document
        """
        print("\n📄 GENERATING WORD DOCUMENT - INDIVIDUAL")
        print("-"*80)
        
        # Professional Background
        professional_bg = intelligence.get("professional_background", "")
        print(f"✓ Professional Background: {len(professional_bg)} chars")
        if not professional_bg or professional_bg.strip() == "":
            print("  ⚠️  EMPTY - Using fallback")
            professional_bg = "No information available."
        else:
            print(f"  Preview: {professional_bg[:200]}...")
        
        doc.add_heading("1. Professional Background", 1)
        # Clean up markdown formatting before adding to document
        cleaned_bg = self._clean_markdown(professional_bg)
        print(f"  📝 Writing to document: {len(cleaned_bg)} chars")
        print(f"  Preview of what's being written: {cleaned_bg[:300]}...")
        doc.add_paragraph(cleaned_bg)
        
        # Education
        education = intelligence.get("education", "")
        print(f"✓ Education: {len(education)} chars")
        if not education or education.strip() == "":
            print("  ⚠️  EMPTY - Using fallback")
            education = "No information available."
        doc.add_heading("2. Education", 1)
        cleaned_education = self._clean_markdown(education)
        print(f"  📝 Writing to document: {len(cleaned_education)} chars")
        doc.add_paragraph(cleaned_education)
        
        # Company Information
        company_info = intelligence.get("company_information", "")
        print(f"\n✓ Company Information: {len(company_info)} chars")
        if not company_info or company_info.strip() == "":
            print("  ⚠️  EMPTY - Using fallback")
            cleaned_company = "No information available."
        else:
            cleaned_company = self._clean_markdown(company_info)
            print(f"  📝 Writing to document: {len(cleaned_company)} chars")
            print(f"  Preview: {cleaned_company[:300]}...")
        doc.add_heading("3. Company Information", 1)
        doc.add_paragraph(cleaned_company)
        
        # Public Presence
        public_presence = intelligence.get("public_presence", "")
        print(f"\n✓ Public Presence: {len(public_presence)} chars")
        if not public_presence or public_presence.strip() == "":
            print("  ⚠️  EMPTY - Using fallback")
            cleaned_public = "No information available."
        else:
            cleaned_public = self._clean_markdown(public_presence)
            print(f"  📝 Writing to document: {len(cleaned_public)} chars")
        doc.add_heading("4. Public Presence", 1)
        doc.add_paragraph(cleaned_public)
        
        # Professional Network
        professional_network = intelligence.get("professional_network", "")
        print(f"\n✓ Professional Network: {len(professional_network)} chars")
        if not professional_network or professional_network.strip() == "":
            print("  ⚠️  EMPTY - Using fallback")
            cleaned_network = "No information available."
        else:
            cleaned_network = self._clean_markdown(professional_network)
            print(f"  📝 Writing to document: {len(cleaned_network)} chars")
        doc.add_heading("5. Professional Network", 1)
        doc.add_paragraph(cleaned_network)
        
        # Recent Activity
        recent_activity = intelligence.get("recent_activity", "")
        print(f"\n✓ Recent Activity: {len(recent_activity)} chars")
        if not recent_activity or recent_activity.strip() == "":
            print("  ⚠️  EMPTY - Using fallback")
            cleaned_activity = "No information available."
        else:
            cleaned_activity = self._clean_markdown(recent_activity)
            print(f"  📝 Writing to document: {len(cleaned_activity)} chars")
        doc.add_heading("6. Recent Activity", 1)
        doc.add_paragraph(cleaned_activity)
        
        print("-"*80)
        
        # Sources
        if intelligence.get("sources"):
            doc.add_heading("Sources", 1)
            for source in intelligence.get("sources", []):
                doc.add_paragraph(source, style="List Bullet")
    
    def _add_organization_content(self, doc: Document, intelligence: Dict[str, Any]):
        """
        Add organization research content to Word document
        """
        # Company Background
        doc.add_heading("1. Company Background", 1)
        doc.add_paragraph(intelligence.get("company_background", "No information available."))
        
        # Leadership Intelligence
        doc.add_heading("2. Leadership Intelligence & Executive Communications", 1)
        leadership = intelligence.get("leadership_intelligence", {})
        if isinstance(leadership, dict):
            doc.add_paragraph(leadership.get("summary", "No information available."))
            
            # Key Executives
            if leadership.get("key_executives"):
                doc.add_heading("Key Executives", 2)
                for exec in leadership.get("key_executives", []):
                    doc.add_paragraph(f"{exec.get('role', '')}: {exec.get('name', '')}", style="List Bullet")
            
            # Executive Messages
            executive_messages = leadership.get("executive_messages", {})
            if executive_messages and any(msg != "No message found" for msg in executive_messages.values()):
                doc.add_heading("Executive Messages & Communications", 2)
                
                for role, message in executive_messages.items():
                    if message and message != "No message found":
                        role_title = role.replace("_", " ").title()
                        para = doc.add_paragraph()
                        run = para.add_run(f"{role_title}: ")
                        run.bold = True
                        para.add_run(message)
            
            # Leadership Targets & Focus Points
            if leadership.get("leadership_targets"):
                doc.add_heading("Leadership Targets & Focus Points", 2)
                for target in leadership.get("leadership_targets", []):
                    doc.add_paragraph(target, style="List Bullet")
        else:
            doc.add_paragraph(str(leadership))
        
        # Financial Information
        doc.add_heading("3. Financial Information (Last 7-10 Years)", 1)
        financial = intelligence.get("financial_information", {})
        if isinstance(financial, dict):
            doc.add_paragraph(financial.get("summary", "No information available."))
        else:
            doc.add_paragraph(str(financial))
        
        # News Intelligence
        doc.add_heading("4. News Intelligence (Last 7-10 Years)", 1)
        news = intelligence.get("news_intelligence", {})
        if isinstance(news, dict):
            # Positive News
            doc.add_heading("Positive News", 2)
            if news.get("positive_news"):
                for item in news.get("positive_news", []):
                    para = doc.add_paragraph(f"[{item.get('date', 'Unknown')}] {item.get('summary', '')}", style="List Bullet")
            else:
                doc.add_paragraph("No positive news items identified.")
            
            # Negative News
            doc.add_heading("Negative News", 2)
            if news.get("negative_news"):
                for item in news.get("negative_news", []):
                    para = doc.add_paragraph(f"[{item.get('date', 'Unknown')}] {item.get('summary', '')}", style="List Bullet")
                    # Highlight negative news
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(200, 0, 0)
            else:
                doc.add_paragraph("No negative news items identified.")
        else:
            doc.add_paragraph(str(news))
        
        # Key Challenges & Risk Intelligence
        doc.add_heading("5. Key Challenges & Public Controversies (Past 10 Years)", 1)
        challenges = intelligence.get("challenges_risks", {})
        if isinstance(challenges, dict):
            doc.add_paragraph(challenges.get("summary", "No information available."))
            
            # External Challenges
            if challenges.get("external_challenges"):
                doc.add_heading("External Challenges", 2)
                for challenge in challenges.get("external_challenges", []):
                    doc.add_paragraph(challenge, style="List Bullet")
            
            # Internal Challenges
            if challenges.get("internal_challenges"):
                doc.add_heading("Internal Challenges", 2)
                for challenge in challenges.get("internal_challenges", []):
                    doc.add_paragraph(challenge, style="List Bullet")
            
            # Public Controversies
            if challenges.get("public_controversies"):
                doc.add_heading("Public Controversies", 2)
                for controversy in challenges.get("public_controversies", []):
                    para = doc.add_paragraph(controversy, style="List Bullet")
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(200, 0, 0)
        else:
            doc.add_paragraph(str(challenges))
        
        # Strategic Priorities
        doc.add_heading("6. Strategic Priorities & Vision", 1)
        doc.add_paragraph(intelligence.get("strategic_priorities", "No information available."))
        
        # Sales Talking Points
        doc.add_heading("7. Sales Talking Points Based on Challenges", 1)
        talking_points = self._generate_talking_points(intelligence)
        for point in talking_points:
            doc.add_paragraph(point, style="List Bullet")
        
        # Recent Activity
        doc.add_heading("8. Recent Activity", 1)
        doc.add_paragraph(intelligence.get("recent_activity", "No information available."))
        
        # Additional Sources Information
        if intelligence.get("additional_sources_info"):
            doc.add_heading("9. Information from Additional Sources", 1)
            doc.add_paragraph(intelligence.get("additional_sources_info", "No information available."))
        
        # User-Provided Sources
        if intelligence.get("additional_sources"):
            doc.add_heading("10. User-Provided Sources", 1)
            for source in intelligence.get("additional_sources", []):
                para = doc.add_paragraph()
                run = para.add_run(f"{source.get('name', 'Unknown')}: ")
                run.bold = True
                para.add_run(source.get('link', 'No link provided'))
        
        # Research Sources
        if intelligence.get("sources"):
            doc.add_heading("Research Sources", 1)
            for source in intelligence.get("sources", []):
                doc.add_paragraph(source, style="List Bullet")
    
    def _generate_talking_points(self, intelligence: Dict[str, Any]) -> list:
        """
        Generate sales talking points based on challenges
        """
        points = []
        challenges = intelligence.get("challenges_risks", {})
        
        if isinstance(challenges, dict):
            # External challenges
            if challenges.get("external_challenges"):
                points.append("Address how our solutions can help navigate external market pressures and competitive challenges.")
            
            # Internal challenges
            if challenges.get("internal_challenges"):
                points.append("Position our offerings as solutions to internal operational and efficiency challenges.")
            
            # Public controversies
            if challenges.get("public_controversies"):
                points.append("Emphasize our commitment to compliance, transparency, and ethical business practices.")
                points.append("Highlight how our solutions can help rebuild trust and improve public perception.")
            
            # Financial stress
            financial = intelligence.get("financial_information", {})
            if isinstance(financial, dict) and financial.get("stress_events"):
                points.append("Focus on ROI, cost reduction, and efficiency improvements.")
                points.append("Emphasize flexible pricing and value-driven solutions.")
        
        # Strategic priorities
        priorities = intelligence.get("strategic_priorities", "")
        if "digital transformation" in str(priorities).lower():
            points.append("Align with their digital transformation initiatives.")
        if "cost reduction" in str(priorities).lower():
            points.append("Emphasize cost savings and operational efficiency.")
        if "compliance" in str(priorities).lower():
            points.append("Highlight compliance and regulatory features.")
        
        if not points:
            points.append("Focus on value proposition and alignment with their strategic goals.")
            points.append("Emphasize partnership and long-term relationship building.")
        
        return points
    
    def _generate_powerpoint(
        self,
        intelligence: Dict[str, Any],
        query: str,
        research_type: Literal["individual", "organization"]
    ) -> str:
        """
        Generate PowerPoint presentation
        """
        prs = Presentation()
        prs.slide_width = PptInches(10)
        prs.slide_height = PptInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Sales Enablement Intelligence Capsule"
        subtitle.text = f"{research_type.title()}: {query}\nGenerated: {datetime.now().strftime('%B %d, %Y')}"
        
        if research_type == "individual":
            self._add_individual_slides(prs, intelligence)
        else:
            self._add_organization_slides(prs, intelligence)
        
        # Save presentation
        filename = f"{query.replace(' ', '_')}_{research_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = f"output/{filename}"
        prs.save(filepath)
        
        return filepath
    
    def _add_individual_slides(self, prs: Presentation, intelligence: Dict[str, Any]):
        """
        Add slides for individual research
        """
        sections = [
            ("Professional Background", intelligence.get("professional_background", "")),
            ("Education", intelligence.get("education", "")),
            ("Company Information", intelligence.get("company_information", "")),
            ("Public Presence", intelligence.get("public_presence", "")),
            ("Recent Activity", intelligence.get("recent_activity", ""))
        ]
        
        for title, content in sections:
            if content:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = slide.shapes.title
                content_shape = slide.placeholders[1]
                
                title_shape.text = title
                content_shape.text = content[:1000]  # Limit content length
    
    def _add_organization_slides(self, prs: Presentation, intelligence: Dict[str, Any]):
        """
        Add slides for organization research
        """
        # Company Background
        self._add_text_slide(prs, "Company Background", intelligence.get("company_background", ""))
        
        # Leadership Intelligence
        leadership = intelligence.get("leadership_intelligence", {})
        if isinstance(leadership, dict):
            content = leadership.get("summary", "")
            if leadership.get("key_executives"):
                content += "\n\nKey Executives:\n"
                for exec in leadership.get("key_executives", []):
                    content += f"• {exec.get('role', '')}: {exec.get('name', '')}\n"
            self._add_text_slide(prs, "Leadership Intelligence", content)
            
            # Executive Messages slide
            executive_messages = leadership.get("executive_messages", {})
            if executive_messages and any(msg != "No message found" for msg in executive_messages.values()):
                exec_content = "Key messages from leadership:\n\n"
                for role, message in executive_messages.items():
                    if message and message != "No message found":
                        role_title = role.replace("_", " ").title()
                        exec_content += f"{role_title}:\n{message}\n\n"
                self._add_text_slide(prs, "Executive Messages & Communications", exec_content)
            
            # Leadership Targets slide
            if leadership.get("leadership_targets"):
                targets_content = "\n".join([
                    f"• {target}" for target in leadership.get("leadership_targets", [])[:5]
                ])
                self._add_text_slide(prs, "Leadership Targets & Focus Points", targets_content)
        
        # Financial Information
        financial = intelligence.get("financial_information", {})
        if isinstance(financial, dict):
            self._add_text_slide(prs, "Financial Information", financial.get("summary", ""))
        
        # News Intelligence
        news = intelligence.get("news_intelligence", {})
        if isinstance(news, dict):
            # Positive News slide
            positive_content = "\n".join([
                f"• [{item.get('date', 'Unknown')}] {item.get('summary', '')}"
                for item in news.get("positive_news", [])[:5]
            ]) or "No positive news items identified."
            self._add_text_slide(prs, "Positive News (Last 7-10 Years)", positive_content)
            
            # Negative News slide
            negative_content = "\n".join([
                f"• [{item.get('date', 'Unknown')}] {item.get('summary', '')}"
                for item in news.get("negative_news", [])[:5]
            ]) or "No negative news items identified."
            self._add_text_slide(prs, "Negative News (Last 7-10 Years)", negative_content)
        
        # Key Challenges & Controversies
        challenges = intelligence.get("challenges_risks", {})
        if isinstance(challenges, dict):
            content = challenges.get("summary", "")
            if challenges.get("external_challenges"):
                content += "\n\nExternal Challenges:\n"
                for ch in challenges.get("external_challenges", [])[:3]:
                    content += f"• {ch}\n"
            if challenges.get("public_controversies"):
                content += "\n\nPublic Controversies:\n"
                for ch in challenges.get("public_controversies", [])[:3]:
                    content += f"• {ch}\n"
            self._add_text_slide(prs, "Key Challenges & Public Controversies", content)
        
        # Strategic Priorities
        self._add_text_slide(prs, "Strategic Priorities", intelligence.get("strategic_priorities", ""))
        
        # Sales Talking Points
        talking_points = self._generate_talking_points(intelligence)
        content = "\n".join([f"• {point}" for point in talking_points])
        self._add_text_slide(prs, "Sales Talking Points Based on Challenges", content)
        
        # Recent Activity
        self._add_text_slide(prs, "Recent Activity", intelligence.get("recent_activity", ""))
        
        # Additional Sources Information
        if intelligence.get("additional_sources_info"):
            self._add_text_slide(prs, "Information from Additional Sources", intelligence.get("additional_sources_info", ""))
        
        # User-Provided Sources
        if intelligence.get("additional_sources"):
            sources_content = "\n".join([
                f"• {source.get('name', 'Unknown')}: {source.get('link', 'No link provided')}"
                for source in intelligence.get("additional_sources", [])
            ])
            self._add_text_slide(prs, "User-Provided Sources", sources_content)
    
    def _add_text_slide(self, prs: Presentation, title: str, content: str):
        """
        Add a text slide with title and content
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        content_shape.text = content[:2000]  # Limit content length
    
    def _clean_markdown(self, text: str) -> str:
        """
        Clean markdown formatting from text for Word document
        """
        if not text:
            return text
        
        # Remove markdown links but keep the text and URL
        text = re.sub(r'\[([^\]]+)\]\(([^\)]+)\)', r'\1 (\2)', text)
        
        # Remove bold/italic markers but keep the text
        text = re.sub(r'\*\*([^\*]+)\*\*', r'\1', text)
        text = re.sub(r'\*([^\*]+)\*', r'\1', text)
        
        # Remove code blocks
        text = re.sub(r'```[^`]*```', '', text, flags=re.DOTALL)
        text = re.sub(r'`([^`]+)`', r'\1', text)
        
        # Clean up multiple spaces
        text = re.sub(r' +', ' ', text)
        
        # Clean up multiple newlines
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        return text.strip()
    
    def _convert_word_to_pdf(self, word_path: str, pdf_path: str):
        """
        Convert Word document to PDF using LibreOffice or pandoc
        """
        try:
            # Try using LibreOffice (common on Linux/Mac)
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", "output", word_path],
                check=True,
                capture_output=True
            )
        except (subprocess.CalledProcessError, FileNotFoundError):
            try:
                # Try using pandoc
                import pypandoc
                pypandoc.convert_file(word_path, "pdf", outputfile=pdf_path)
            except Exception:
                # If conversion fails, just return the word path
                # User can manually convert or we can add another method
                pass

